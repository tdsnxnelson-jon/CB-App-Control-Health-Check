"""Generate a PPTX that compares four severity-tile color directions.

Usage:
    python tools/tile_palette_preview.py [output-path]
"""
import os
import sys
from pathlib import Path

from pptx.dml.color import RGBColor

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from healthcheck.report import pptx_helpers as ph


PALETTES = [
    (
        "Professional Contrast",
        "Deep, legible status colors for a report that still feels restrained.",
        {
            "critical": ("A61B1B", "F6DEDE"),
            "warning": ("B45309", "FCE8CF"),
            "caution": ("8A6500", "F8F0CE"),
            "ok": ("176B3A", "DDEFE3"),
            "info": ("176B87", "DCEBF1"),
        },
    ),
    (
        "Dark Executive",
        "Charcoal tiles with white type and decisive severity accents.",
        {
            "critical": ("FF6B6B", "352126"),
            "warning": ("FFB454", "362B1F"),
            "caution": ("F4D35E", "35321F"),
            "ok": ("63D297", "1F332B"),
            "info": ("69C6E8", "1D3038"),
        },
    ),
    (
        "Brighter Modern",
        "Clear, saturated colors with a clean contemporary feel.",
        {
            "critical": ("D92D20", "FFE2DE"),
            "warning": ("E56A00", "FFE9D2"),
            "caution": ("A87900", "FFF1BF"),
            "ok": ("07804A", "D7F5E5"),
            "info": ("087EA4", "D7F3FA"),
        },
    ),
    (
        "Minimal Neutral",
        "Graphite tiles where color is reserved for the status rail and label.",
        {
            "critical": ("D92D20", "F1F3F5"),
            "warning": ("C56A00", "F1F3F5"),
            "caution": ("8A6500", "F1F3F5"),
            "ok": ("16804A", "F1F3F5"),
            "info": ("087EA4", "F1F3F5"),
        },
    ),
]


def _rgb(hex_value: str) -> RGBColor:
    return RGBColor(*bytes.fromhex(hex_value))


def _apply_palette(values) -> None:
    for severity, (accent, tone) in values.items():
        ph.SEVERITY_COLORS[severity] = _rgb(accent)
        ph.SEVERITY_TONES[severity] = _rgb(tone)


def build_preview(output_path: str) -> str:
    prs = ph.create_deck()
    findings = [
        ("critical", "23 endpoints have not checked in for more than 7 days."),
        ("warning", "Approval activity is above the normal daily volume."),
        ("caution", "Five custom rules use broad wildcard paths."),
        ("ok", "Agent sync coverage is meeting the reporting target."),
    ]

    for title, subtitle, palette in PALETTES:
        _apply_palette(palette)
        slide = ph.add_content_slide(prs, title)
        ph._add_text(slide, ph.MARGIN, 1.12, ph.CONTENT_W, 0.3, subtitle, size=12, color=ph.MUTED_TEXT)
        ph.add_metric_strip(
            slide,
            [("Critical", "2", "critical"), ("Warning", "7", "warning"), ("Passing", "38", "ok")],
            top=1.55,
        )
        ph.add_finding_cards(slide, findings, top=2.52, height=4.35, show_recommendations=False)

        if title == "Dark Executive":
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.color.rgb = ph.WHITE
                        for run in paragraph.runs:
                            run.font.color.rgb = ph.WHITE
            for paragraph in slide.shapes[2].text_frame.paragraphs:
                paragraph.font.color.rgb = ph.NAVY
                for run in paragraph.runs:
                    run.font.color.rgb = ph.NAVY
            for paragraph in slide.shapes[3].text_frame.paragraphs:
                paragraph.font.color.rgb = ph.MUTED_TEXT
                for run in paragraph.runs:
                    run.font.color.rgb = ph.MUTED_TEXT

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    default_output = os.path.join(os.path.dirname(__file__), "tile_palette_preview.pptx")
    print(build_preview(sys.argv[1] if len(sys.argv) > 1 else default_output))