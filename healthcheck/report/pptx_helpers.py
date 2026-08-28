"""
Generic python-pptx helpers for building slides.

Everything is drawn on the "Blank" layout with explicit positions instead
of relying on a template's placeholder geometry - this guarantees content
never overlaps the title regardless of which default/looked-up template
python-pptx ends up using, and gives predictable, consistent spacing.
"""
from typing import Iterable, List, Sequence, Tuple
import math

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- canvas geometry (16:9) ----------------------------------------------
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.45
TITLE_TOP = 0.3
TITLE_H = 0.7
CONTENT_TOP = TITLE_TOP + TITLE_H + 0.2
CONTENT_W = SLIDE_W - 2 * MARGIN
CONTENT_H = SLIDE_H - CONTENT_TOP - MARGIN

BLANK_LAYOUT = 6

NAVY = RGBColor(0x1F, 0x2D, 0x3D)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2B, 0x2B, 0x2B)
LIGHT_GRID = RGBColor(0xE8, 0xEC, 0xF1)
SOFT_BG = RGBColor(0xF8, 0xFA, 0xFC)
MUTED_TEXT = RGBColor(0x63, 0x6A, 0x73)

SEVERITY_COLORS = {
    "critical": RGBColor(0xC0, 0x00, 0x00),
    "warning": RGBColor(0xE3, 0x7B, 0x0F),
    "caution": RGBColor(0xB8, 0x8A, 0x00),
    "ok": RGBColor(0x1E, 0x7B, 0x34),
    "info": RGBColor(0x2E, 0x75, 0xB6),
}
SEVERITY_LABELS = {"critical": "CRITICAL", "warning": "WARNING", "caution": "CAUTION", "ok": "OK", "info": "INFO"}
SEVERITY_TONES = {
    "critical": RGBColor(0xFB, 0xE9, 0xE9),
    "warning": RGBColor(0xFE, 0xF1, 0xE5),
    "caution": RGBColor(0xFB, 0xF5, 0xDD),
    "ok": RGBColor(0xE8, 0xF4, 0xEC),
    "info": RGBColor(0xE9, 0xF1, 0xFA),
}
PIE_COLORS = [
    RGBColor(0x4F, 0x81, 0xBD),
    RGBColor(0xC0, 0x50, 0x4D),
    RGBColor(0x9B, 0xBB, 0x59),
    RGBColor(0x80, 0x64, 0xA2),
    RGBColor(0x4B, 0xAC, 0xC6),
    RGBColor(0xF7, 0x96, 0x46),
]


def create_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_text(slide, left, top, width, height, text, size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_rich_bullets(slide, items: Sequence, left=MARGIN, top=CONTENT_TOP, width=CONTENT_W, height=CONTENT_H, font_size=15):
    """Renders nested bullets from (level, runs) items.

    Each run is (text, bold, italic) or (text, bold, italic, hyperlink_url).
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, item in enumerate(items):
        level, runs = item
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.space_after = Pt(4 if level == 0 else 2)
        p.line_spacing = 1.05
        prefix = "▪  " if level == 0 else "-  "
        prefix_run = p.add_run()
        prefix_run.text = ("    " * level) + prefix
        prefix_run.font.size = Pt(font_size - 1)
        prefix_run.font.color.rgb = ACCENT
        for run_data in runs:
            text, bold, italic = run_data[:3]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = DARK_TEXT
            if len(run_data) > 3 and run_data[3]:
                run.font.color.rgb = ACCENT
                run.font.underline = True
                run.hyperlink.address = run_data[3]
    return box


def _set_shape_text(shape, text, size=12, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def add_metric_strip(slide, metrics: Sequence[Tuple[str, str, str]], left=MARGIN, top=CONTENT_TOP, width=CONTENT_W, height=0.72):
    """Draws compact executive-summary metric chips.

    metrics are (label, value, severity_or_None) tuples.
    """
    if not metrics:
        return
    gap = 0.16
    card_w = (width - gap * (len(metrics) - 1)) / len(metrics)
    for i, (label, value, severity) in enumerate(metrics):
        x = left + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(card_w), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_GRID
        if severity:
            _add_rect(slide, x, top, 0.06, height, SEVERITY_COLORS.get(severity, ACCENT))
        # longer values (e.g. "77 (C)") get a wider box and smaller font so they stay on one line
        value_text = str(value)
        value_w = 0.55 if len(value_text) <= 2 else min(1.05, 0.55 + 0.11 * (len(value_text) - 2))
        value_size = 18 if len(value_text) <= 2 else 14
        value_box = _add_text(slide, x + 0.16, top + 0.12, value_w, 0.36, value_text, size=value_size, bold=True, color=SEVERITY_COLORS.get(severity, NAVY), anchor=MSO_ANCHOR.MIDDLE)
        value_box.text_frame.word_wrap = False
        label_x = x + 0.16 + value_w + 0.07
        _add_text(slide, label_x, top + 0.16, x + card_w - 0.16 - label_x, 0.32, label, size=10, bold=True, color=MUTED_TEXT, anchor=MSO_ANCHOR.MIDDLE)


def _normalize_findings(items: Sequence, max_items=None) -> List[Tuple]:
    source = items[:max_items] if max_items else items
    normalized: List[Tuple] = []
    for item in source:
        if isinstance(item, tuple):
            severity, message = item[0], item[1]
            recommendation = item[2] if len(item) > 2 else None
            normalized.append((severity, message, recommendation))
        else:
            normalized.append((None, item, None))
    return normalized


def _estimate_text_height(text: str, width_in: float, font_size: int, bold: bool = False, line_spacing: float = 1.15) -> float:
    """Rough estimate (inches) of wrapped text block height, used to size
    finding cards so long messages don't overlap the recommendation subtext."""
    if not text:
        return 0.0
    char_w_in = font_size / 72.0 * (0.62 if bold else 0.55)
    chars_per_line = max(1, int(width_in / char_w_in))
    lines = max(1, math.ceil(len(text) / chars_per_line))
    return lines * (font_size / 72.0 * line_spacing)


def add_finding_cards(slide, items: Sequence, left=MARGIN, top=CONTENT_TOP, width=CONTENT_W, height=CONTENT_H, show_recommendations=True, max_items=8, compact=False):
    """Renders finding text as a card grid instead of a dense bullet list.

    items may be (severity, message), (severity, message, recommendation), or
    plain strings. Recommendation text is optional and can be hidden for
    executive-facing slides.
    """
    normalized = _normalize_findings(items, max_items)

    if not normalized:
        normalized = [("ok", "No notable findings in the analyzed data.", None)]

    count = len(normalized)
    cols = 3 if compact and count >= 5 else 2 if count > 2 else 1
    rows = (count + cols - 1) // cols
    gap_x = 0.18
    gap_y = 0.16
    card_w = (width - gap_x * (cols - 1)) / cols
    default_card_h = (height - gap_y * (rows - 1)) / rows
    if compact:
        default_card_h = min(default_card_h, 1.55)

    message_size = 11 if compact else 12 if count <= 6 else 10
    rec_size = 7 if compact else 8
    text_w = card_w - 0.42
    header_h = 0.42

    # Size each row to fit its tallest card's estimated content instead of a
    # fixed split between message/recommendation, so long text doesn't overlap.
    row_heights = []
    for row in range(rows):
        needed = default_card_h
        for col in range(cols):
            i = row * cols + col
            if i >= count:
                continue
            _, message, recommendation = normalized[i]
            message_h = _estimate_text_height(message, text_w, message_size, bold=True)
            rec_text = f"Next: {recommendation}" if show_recommendations and recommendation else ""
            rec_h = _estimate_text_height(rec_text, text_w, rec_size) if rec_text else 0.0
            gap = 0.1 if rec_h else 0.0
            needed = max(needed, header_h + message_h + gap + rec_h + 0.12)
        row_heights.append(needed)

    for i, (severity, message, recommendation) in enumerate(normalized):
        col = i % cols
        row = i // cols
        card_h = row_heights[row]
        x = left + col * (card_w + gap_x)
        y = top + sum(row_heights[:row]) + row * gap_y
        severity_color = SEVERITY_COLORS.get(severity, ACCENT)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_GRID
        _add_rect(slide, x, y, 0.08, card_h, severity_color)

        label = SEVERITY_LABELS.get(severity, "NOTE")
        _add_text(slide, x + 0.22, y + 0.12, 1.2, 0.22, label, size=8, bold=True, color=severity_color, anchor=MSO_ANCHOR.MIDDLE)

        message_h = _estimate_text_height(message, text_w, message_size, bold=True)
        rec_text = f"Next: {recommendation}" if show_recommendations and recommendation else ""
        rec_h = _estimate_text_height(rec_text, text_w, rec_size) if rec_text else 0.0
        message_box_h = max(0.38, card_h - header_h - (0.1 + rec_h if rec_h else 0.0) - 0.12)
        _add_text(slide, x + 0.22, y + header_h, text_w, message_box_h, message, size=message_size, bold=True, color=DARK_TEXT)

        if rec_text:
            _add_text(slide, x + 0.22, y + card_h - rec_h - 0.1, text_w, rec_h, rec_text, size=rec_size, color=MUTED_TEXT)


def add_findings_dashboard(slide, items: Sequence, left=MARGIN, top=CONTENT_TOP, width=CONTENT_W, height=CONTENT_H, max_cards=6):
    normalized = _normalize_findings(items)
    if not normalized:
        add_finding_cards(slide, [("ok", "No notable findings in the analyzed data.")], left, top, width, height)
        return

    order = {"critical": 0, "warning": 1, "caution": 2, "info": 3, "ok": 4, None: 5}
    normalized.sort(key=lambda item: order.get(item[0], 5))
    counts = {sev: sum(1 for item in normalized if item[0] == sev) for sev in ("critical", "warning", "caution", "ok")}
    metrics = [
        ("Critical", counts["critical"], "critical" if counts["critical"] else "ok"),
        ("Warning", counts["warning"], "warning" if counts["warning"] else "ok"),
        ("Caution", counts["caution"], "caution" if counts["caution"] else "ok"),
        ("Passing", counts["ok"], "ok"),
    ]
    add_metric_strip(slide, metrics, left=left, top=top, width=width, height=0.64)
    add_finding_cards(slide, normalized[:max_cards], left=left, top=top + 0.88, width=width, height=height - 0.88, show_recommendations=True, max_items=max_cards, compact=True)
    remaining = len(normalized) - max_cards
    if remaining > 0:
        _add_text(slide, left, top + height - 0.22, width, 0.18, f"+ {remaining} additional low-priority item(s) not shown on this slide.", size=8, color=MUTED_TEXT, align=PP_ALIGN.RIGHT)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _add_rect(slide, 0, SLIDE_H / 2 + 0.4, SLIDE_W, 0.06, ACCENT)
    _add_text(slide, 1, SLIDE_H / 2 - 1.3, SLIDE_W - 2, 1.1, title, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(slide, 1, SLIDE_H / 2 + 0.55, SLIDE_W - 2, 1.0, subtitle, size=18, color=RGBColor(0xC9, 0xD6, 0xE3), align=PP_ALIGN.CENTER)


def add_section_slide(prs: Presentation, title: str) -> None:
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, ACCENT)
    _add_rect(slide, 1, SLIDE_H / 2 - 0.05, 1.4, 0.08, WHITE)
    _add_text(slide, 1, SLIDE_H / 2 - 1.0, SLIDE_W - 2, 0.9, title, size=32, bold=True, color=WHITE)


def add_content_slide(prs: Presentation, title: str):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, SLIDE_W, TITLE_TOP + TITLE_H + 0.08, RGBColor(0xF2, 0xF4, 0xF7))
    _add_rect(slide, 0, TITLE_TOP + TITLE_H + 0.08, SLIDE_W, 0.04, ACCENT)
    _add_text(slide, MARGIN, TITLE_TOP, SLIDE_W - 2 * MARGIN, TITLE_H, title, size=24, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def add_footnote(slide, text: str):
    _add_text(slide, MARGIN, SLIDE_H - 0.35, CONTENT_W, 0.18, text, size=8, color=MUTED_TEXT)


def add_bullets(slide, items: Sequence, left=MARGIN, top=CONTENT_TOP, width=CONTENT_W, height=CONTENT_H, font_size=15):
    """items: sequence of (severity_or_None, message) or (severity, message,
    recommendation) tuples, or plain strings. A recommendation (when given,
    typically for anything above INFO) renders as an indented sub-line."""
    normalized: List[Tuple] = []
    for item in items:
        if isinstance(item, tuple):
            severity, message = item[0], item[1]
            recommendation = item[2] if len(item) > 2 else None
            normalized.append((severity, message, recommendation))
        else:
            normalized.append((None, item, None))

    # shrink font if there are a lot of *rendered lines* (a recommendation
    # sub-line effectively doubles an item's height) so nothing overflows
    effective_lines = sum(2 if rec else 1 for _, _, rec in normalized) or 1
    fitted_size = font_size if effective_lines <= 10 else max(9, font_size - (effective_lines - 10) // 2)
    rec_size = max(8, fitted_size - 2)

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for severity, message, recommendation in normalized:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(2 if recommendation else 6)
        p.line_spacing = 1.1
        if severity and severity in SEVERITY_COLORS:
            tag = p.add_run()
            tag.text = f"{SEVERITY_LABELS[severity]}  "
            tag.font.size = Pt(fitted_size)
            tag.font.bold = True
            tag.font.color.rgb = SEVERITY_COLORS[severity]
        else:
            bullet = p.add_run()
            bullet.text = "\u25aa  "
            bullet.font.size = Pt(fitted_size)
            bullet.font.color.rgb = ACCENT
        body = p.add_run()
        body.text = str(message)
        body.font.size = Pt(fitted_size)
        body.font.color.rgb = DARK_TEXT
        if recommendation:
            rec_p = tf.add_paragraph()
            rec_p.space_after = Pt(8)
            rec_p.line_spacing = 1.05
            rec_run = rec_p.add_run()
            rec_run.text = f"\u2192 Recommendation: {recommendation}"
            rec_run.font.size = Pt(rec_size)
            rec_run.font.italic = True
            rec_run.font.color.rgb = RGBColor(0x55, 0x5A, 0x60)


def _text_width_estimate(text: str, size_pt: float) -> float:
    return len(str(text)) * size_pt * 0.0092  # rough inches-per-char at given point size


def _truncate_to_width(text, size_pt: float, width_in: float) -> str:
    """Prevents word-wrap from silently growing a table row taller than
    planned - the actual cause of tables overflowing the slide - by
    hard-truncating cell text that wouldn't fit its column on one line."""
    text = "" if text is None else str(text)
    max_chars = max(4, int(width_in / (size_pt * 0.0092)))
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "\u2026"


ROWS_PER_PAGE_BY_FONT = {14: 13, 12: 15, 11: 17, 10: 19, 9: 21, 8: 23, 7: 25}


def _rows_per_page(font_size: int) -> int:
    return ROWS_PER_PAGE_BY_FONT.get(font_size, 15)


def _auto_col_widths(rows: List[Sequence], total_width: float, font_size: int) -> List[float]:
    n_cols = len(rows[0])
    widths = [0.0] * n_cols
    for row in rows:
        for c, val in enumerate(row):
            widths[c] = max(widths[c], _text_width_estimate(val, font_size) + 0.22)
    widths = [max(w, 0.55) for w in widths]

    # Let compact tables stay compact instead of stretching every column to
    # the full content area. If the natural width is too wide, scale it down
    # to the available space and rely on _truncate_to_width to avoid wrapping.
    cap = total_width * (0.55 if n_cols > 2 else 0.85)
    widths = [min(w, cap) for w in widths]
    total = sum(widths)
    if total <= total_width:
        return widths
    scale = total_width / total
    return [max(0.45, w * scale) for w in widths]


def _fit_col_widths(widths: Sequence[float], max_width: float) -> List[float]:
    widths = list(widths)
    total = sum(widths)
    if total <= max_width:
        return widths
    scale = max_width / total
    return [max(0.45, w * scale) for w in widths]


def add_table(
    slide,
    rows: List[Sequence],
    left=MARGIN,
    top=CONTENT_TOP,
    width=CONTENT_W,
    height=CONTENT_H,
    font_size=12,
    col_widths: Sequence[float] = None,
    header=True,
    max_rows=30,
    center=False,
):
    """rows[0] is treated as the header row when header=True. Column widths
    auto-size to content (falls back to even split) and font size/row
    count shrink automatically so the table never overflows the slide."""
    if len(rows) - 1 > max_rows:
        kept = rows[:1] + rows[1:max_rows + 1]
        kept.append([f"+ {len(rows) - 1 - max_rows} more row(s) not shown"] + [""] * (len(rows[0]) - 1))
        rows = kept

    n_rows = len(rows)
    min_row_h = 0.28
    size = font_size
    while n_rows * min_row_h > height and size > 7:
        size -= 1
        min_row_h = max(0.22, min_row_h - 0.01)
    table_height = min(height, n_rows * max(min_row_h, 0.22))

    widths = _fit_col_widths(col_widths, width) if col_widths else _auto_col_widths(rows, width, size)
    table_width = min(width, sum(widths))
    if center:
        left += (width - table_width) / 2
        top += (height - table_height) / 2
    shape = slide.shapes.add_table(n_rows, len(rows[0]), Inches(left), Inches(top), Inches(table_width), Inches(table_height))
    table = shape.table

    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = _truncate_to_width(val, size, widths[c])
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(size)
                if header and r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_TEXT
            if header and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRID
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table


def add_table_slides(prs, title: str, rows: List[Sequence], font_size=11, col_widths: Sequence[float] = None, header=True, max_total_rows=150):
    """Adds one or more content slides to hold `rows`, splitting into
    additional "(cont'd)" slides instead of letting a long table overflow
    a single slide. Use this instead of add_content_slide+add_table for
    any table whose row count isn't already capped to ~15-20."""
    header_row = rows[0] if header else None
    data_rows = rows[1:] if header else rows

    if len(data_rows) > max_total_rows:
        data_rows = data_rows[:max_total_rows]
        data_rows.append([f"+ {len(rows) - 1 - max_total_rows} more row(s) not shown"] + [""] * (len(rows[0]) - 1))

    rows_per_page = _rows_per_page(font_size)
    chunks = [data_rows[i:i + rows_per_page] for i in range(0, len(data_rows), rows_per_page)] or [[]]
    n_pages = len(chunks)

    slides = []
    for i, chunk in enumerate(chunks):
        page_title = title if n_pages == 1 else f"{title} ({i + 1}/{n_pages})"
        slide = add_content_slide(prs, page_title)
        page_rows = ([header_row] if header else []) + chunk
        add_table(slide, page_rows, font_size=font_size, col_widths=col_widths, header=header, max_rows=len(page_rows), center=True)
        slides.append(slide)
    return slides


def color_cell(table, row: int, col: int, severity: str):
    color = SEVERITY_COLORS.get(severity)
    if not color:
        return
    cell = table.cell(row, col)
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def _rotate_category_labels(chart, degrees=-45):
    cat_ax = chart.category_axis._element
    txPr = cat_ax.find(qn("c:txPr"))
    if txPr is None:
        txPr = cat_ax.makeelement(qn("c:txPr"), {})
        cat_ax.append(txPr)
    bodyPr = txPr.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = txPr.makeelement(qn("a:bodyPr"), {})
        txPr.insert(0, bodyPr)
    bodyPr.set("rot", str(int(degrees * 60000)))


def add_pie_chart(slide, title: str, categories: Iterable[str], values: Iterable[float], left=1.5, top=CONTENT_TOP, width=CONTENT_W - 3.0, height=CONTENT_H, show_legend=True, colors=None):
    categories, values = list(categories), list(values)

    # too many slices makes pie labels unreadable - a horizontal bar reads
    # far better once there are more than ~6 categories.
    if len(categories) > 6:
        return add_bar_chart(slide, title, categories, {"Count": values}, left=MARGIN, top=top, width=width + 3.0, height=height, horizontal=True)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Series 1", values)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(left), Inches(top), Inches(width), Inches(height), chart_data)
    chart = gframe.chart
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False
    plot.data_labels.show_category_name = False
    plot.data_labels.number_format = "0%"
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(10)
    if colors:
        for point, color in zip(chart.series[0].points, colors):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color
    return chart


def add_line_chart(slide, title: str, categories: Iterable[str], series: dict, left=MARGIN, top=CONTENT_TOP, width=CONTENT_W, height=CONTENT_H):
    categories = list(categories)
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series.items():
        chart_data.add_series(name, list(values))
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(left), Inches(top), Inches(width), Inches(height), chart_data)
    chart = gframe.chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.number_format = "#,##0"
    chart.value_axis.tick_labels.number_format_is_linked = False
    if len(categories) > 10:
        _rotate_category_labels(chart, -45)
    return chart


def add_bar_chart(slide, title: str, categories: Iterable[str], series: dict, left=MARGIN, top=CONTENT_TOP, width=CONTENT_W, height=CONTENT_H, horizontal=False):
    categories = list(categories)
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series.items():
        chart_data.add_series(name, list(values))
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    gframe = slide.shapes.add_chart(chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data)
    chart = gframe.chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.number_format = "#,##0"
    chart.value_axis.tick_labels.number_format_is_linked = False
    if not horizontal and len(categories) > 8:
        _rotate_category_labels(chart, -45)
    return chart
